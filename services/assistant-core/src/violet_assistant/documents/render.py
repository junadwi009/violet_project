from __future__ import annotations

import io
from typing import Any


def _as_list(value: Any) -> list:
    if value is None:
        return []
    return value if isinstance(value, list) else [value]


def _text(value: Any) -> str:
    return "" if value is None else str(value)


def render_docx(spec: dict) -> bytes:
    """Render a DOCX from a spec: {title, subtitle?, sections:[{heading, paragraphs?, bullets?, table?}]}."""
    from docx import Document

    document = Document()
    if spec.get("title"):
        document.add_heading(_text(spec["title"]), level=0)
    if spec.get("subtitle"):
        document.add_paragraph(_text(spec["subtitle"])).italic = True

    for section in _as_list(spec.get("sections")):
        if not isinstance(section, dict):
            document.add_paragraph(_text(section))
            continue
        if section.get("heading"):
            document.add_heading(_text(section["heading"]), level=1)
        for paragraph in _as_list(section.get("paragraphs")):
            document.add_paragraph(_text(paragraph))
        for bullet in _as_list(section.get("bullets")):
            document.add_paragraph(_text(bullet), style="List Bullet")
        table_spec = section.get("table")
        if isinstance(table_spec, dict) and table_spec.get("headers"):
            headers = _as_list(table_spec.get("headers"))
            rows = _as_list(table_spec.get("rows"))
            table = document.add_table(rows=1, cols=len(headers))
            try:
                table.style = "Light Grid Accent 1"
            except (KeyError, ValueError):
                pass
            for index, header in enumerate(headers):
                table.rows[0].cells[index].text = _text(header)
            for row in rows:
                cells = table.add_row().cells
                for index, value in enumerate(_as_list(row)[: len(headers)]):
                    cells[index].text = _text(value)

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def render_pptx(spec: dict) -> bytes:
    """Render a PPTX from a spec: {title, subtitle?, slides:[{title, bullets?, notes?}]}."""
    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()

    title_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_layout)
    title_slide.shapes.title.text = _text(spec.get("title") or "Presentation")
    if title_slide.placeholders and len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = _text(spec.get("subtitle"))

    content_layout = presentation.slide_layouts[1]
    for slide_spec in _as_list(spec.get("slides")):
        if not isinstance(slide_spec, dict):
            slide_spec = {"title": _text(slide_spec)}
        slide = presentation.slides.add_slide(content_layout)
        slide.shapes.title.text = _text(slide_spec.get("title"))

        body = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
        bullets = _as_list(slide_spec.get("bullets"))
        if body is not None and bullets:
            body.clear()
            for index, bullet in enumerate(bullets):
                para = body.paragraphs[0] if index == 0 else body.add_paragraph()
                para.text = _text(bullet)
                para.font.size = Pt(18)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = _text(notes)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
